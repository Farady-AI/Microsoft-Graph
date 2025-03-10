from fastapi import FastAPI, HTTPException
import logging
import os
import msal
import requests
from datetime import datetime, timedelta
from pydantic import BaseModel
from pptx import Presentation
from docx import Document
import pandas as pd
from starlette.responses import FileResponse
import uvicorn
import gunicorn
import openai
from dotenv import load_dotenv
from fastapi.middleware.cors import CORSMiddleware
from starlette.middleware.sessions import SessionMiddleware
from pptx.util import Inches
from typing import Optional
from fastapi.responses import RedirectResponse
from fastapi import Request, Response

# Setup logging
logging.basicConfig(level=logging.INFO)

# Define request models
class DocumentRequest(BaseModel):
    document_type: str
    title: str
    content: str
    output_format: str = "docx"  # "docx" or "pptx"

class TextRequest(BaseModel):
    prompt: str
    structured_output: bool = False
    document_request: DocumentRequest = None

# Load environment variables
load_dotenv()
CLIENT_ID = os.getenv("MSCLIENTID")
CLIENT_SECRET = os.getenv("MSCLIENTSECRET")
TENANT_ID = os.getenv("MSTENANTID")
REDIRECT_URI = os.getenv("REDIRECT_URI")  # Must match Azure settings
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")  # OpenAI API Key

# Validate required environment variables
def validate_env_vars():
    missing_vars = []
    if not CLIENT_ID:
        missing_vars.append("MSCLIENTID")
    if not CLIENT_SECRET:
        missing_vars.append("MSCLIENTSECRET")
    if not TENANT_ID:
        missing_vars.append("MSTENANTID")
    if not REDIRECT_URI:
        missing_vars.append("REDIRECT_URI")
    
    if missing_vars:
        raise ValueError(f"Missing required environment variables: {', '.join(missing_vars)}")

    # OpenAI is optional
    if not OPENAI_API_KEY:
        logging.warning("OPENAI_API_KEY is not set. Text generation will use fallback mode.")

# Initialize OpenAI client
def init_openai():
    if not OPENAI_API_KEY:
        return None
    try:
        client = openai.OpenAI()
        return client
    except Exception as e:
        logging.error(f"Failed to initialize OpenAI client: {str(e)}")
        return None

AUTHORITY = f"https://login.microsoftonline.com/{TENANT_ID}"
SCOPE = ["https://graph.microsoft.com/Mail.Send", "https://graph.microsoft.com/User.Read"]

app = FastAPI()

# Add session middleware with a secure key
app.add_middleware(
    SessionMiddleware,
    secret_key=os.urandom(32).hex(),  # Generate a random secret key
    max_age=3600  # Session expires in 1 hour
)

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allows all origins for testing
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

user_tokens = {}  # Dictionary to store user tokens (Use a database in production)

openai.api_key = OPENAI_API_KEY  # Set OpenAI API Key

@app.on_event("startup")
async def startup_event():
    """Validate configuration on startup"""
    try:
        validate_env_vars()
        logging.info("Server configuration validated successfully")
    except Exception as e:
        logging.error(f"Startup validation failed: {str(e)}")
        # We'll continue running but with limited functionality

# Lazy initialization of OpenAI client
_openai_client = None

def get_openai_client():
    global _openai_client
    if _openai_client is None:
        _openai_client = init_openai()
    return _openai_client

# Add these classes after the existing model definitions
class AuthError(Exception):
    def __init__(self, message: str, error_code: Optional[str] = None):
        self.message = message
        self.error_code = error_code
        super().__init__(self.message)

class MSGraphAuth:
    def __init__(self, client_id: str, client_secret: str, tenant_id: str, redirect_uri: str):
        if not all([client_id, client_secret, tenant_id, redirect_uri]):
            raise ValueError("Missing required Microsoft Graph credentials")
        self.client_id = client_id
        self.client_secret = client_secret
        self.tenant_id = tenant_id
        self.redirect_uri = redirect_uri
        self.authority = f"https://login.microsoftonline.com/{tenant_id}"
        self._confidential_client = None
        self._public_client = None

    @property
    def confidential_client(self):
        if not self._confidential_client:
            self._confidential_client = msal.ConfidentialClientApplication(
                self.client_id,
                authority=self.authority,
                client_credential=self.client_secret
            )
        return self._confidential_client

    @property
    def public_client(self):
        if not self._public_client:
            self._public_client = msal.PublicClientApplication(
                self.client_id,
                authority=self.authority
            )
        return self._public_client

# Initialize MSGraphAuth after environment variables are loaded
try:
    graph_auth = MSGraphAuth(CLIENT_ID, CLIENT_SECRET, TENANT_ID, REDIRECT_URI)
except ValueError as e:
    logging.error(f"Failed to initialize Microsoft Graph Auth: {str(e)}")
    graph_auth = None

@app.get("/auth/login")
async def login(request: Request):
    """
    Initiate the Microsoft Graph authentication flow
    """
    try:
        if not graph_auth:
            raise AuthError("Microsoft Graph authentication is not configured")

        # Generate and store state parameter for CSRF protection
        state = os.urandom(16).hex()
        request.session["auth_state"] = state

        auth_url = graph_auth.public_client.get_authorization_request_url(
            SCOPE,
            redirect_uri=REDIRECT_URI,
            state=state
        )
        return {"auth_url": auth_url}
    except AuthError as ae:
        logging.error(f"Authentication configuration error: {str(ae)}")
        raise HTTPException(status_code=503, detail=str(ae))
    except Exception as e:
        logging.error(f"Unexpected error in login: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal server error during authentication")

@app.get("/auth/callback")
async def auth_callback(request: Request, code: str, state: Optional[str] = None):
    """
    Handle the Microsoft Graph authentication callback
    """
    try:
        if not graph_auth:
            raise AuthError("Microsoft Graph authentication is not configured")

        # Verify state parameter if provided
        stored_state = request.session.get("auth_state")
        if state and stored_state and state != stored_state:
            raise AuthError("Invalid state parameter", "invalid_state")

        # Attempt to acquire token
        result = graph_auth.confidential_client.acquire_token_by_authorization_code(
            code,
            scopes=SCOPE,
            redirect_uri=REDIRECT_URI
        )

        if "error" in result:
            raise AuthError(
                f"Failed to acquire token: {result.get('error_description')}",
                result.get('error')
            )

        if "access_token" not in result:
            raise AuthError("No access token received")

        # Get user info
        headers = {"Authorization": f"Bearer {result['access_token']}"}
        user_response = requests.get(
            "https://graph.microsoft.com/v1.0/me",
            headers=headers
        )
        user_response.raise_for_status()
        user_info = user_response.json()

        # Extract user email
        user_email = user_info.get("mail") or user_info.get("userPrincipalName")
        if not user_email:
            raise AuthError("Unable to fetch user email")

        # Store token with expiration
        expiration_time = datetime.utcnow() + timedelta(seconds=result["expires_in"])
        result["expires_at"] = expiration_time.timestamp()
        user_tokens[user_email] = result

        logging.info(f"Authentication successful for user: {user_email}")
        return {
            "message": "Login successful",
            "user": user_email,
            "token_expires_at": expiration_time.isoformat()
        }

    except AuthError as ae:
        logging.error(f"Authentication error: {ae.message}")
        raise HTTPException(
            status_code=400,
            detail={"message": ae.message, "error_code": ae.error_code}
        )
    except requests.exceptions.RequestException as re:
        logging.error(f"Microsoft Graph API error: {str(re)}")
        raise HTTPException(status_code=502, detail="Failed to communicate with Microsoft Graph API")
    except Exception as e:
        logging.error(f"Unexpected error in auth callback: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal server error during authentication")

# Add a helper function to verify tokens
async def verify_token(email: str) -> bool:
    """
    Verify if a user's token is valid and not expired
    """
    if email not in user_tokens:
        return False
    
    token_data = user_tokens[email]
    if datetime.utcnow().timestamp() >= token_data["expires_at"]:
        # Token expired, try to refresh
        try:
            new_token = graph_auth.confidential_client.acquire_token_silent(
                SCOPE,
                account=None
            )
            if new_token:
                expiration_time = datetime.utcnow() + timedelta(seconds=new_token["expires_in"])
                new_token["expires_at"] = expiration_time.timestamp()
                user_tokens[email] = new_token
                return True
        except Exception as e:
            logging.error(f"Token refresh failed for {email}: {str(e)}")
            return False
        
    return True

@app.get("/generate-ppt")
def generate_ppt():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Generated PowerPoint Slide"
    prs.save("generated_ppt.pptx")
    return FileResponse("generated_ppt.pptx")

@app.get("/generate-doc")
def generate_doc():
    doc = Document()
    doc.add_paragraph("Generated Word Document")
    doc.save("generated_doc.docx")
    return FileResponse("generated_doc.docx")

@app.get("/generate-excel")
def generate_excel():
    df = pd.DataFrame({"Column1": ["Data1", "Data2"], "Column2": ["MoreData1", "MoreData2"]})
    df.to_excel("generated_excel.xlsx", index=False)
    return FileResponse("generated_excel.xlsx")

@app.post("/generate-text")
async def generate_text(request: TextRequest):
    """
    Enhanced endpoint that handles both free-form text generation and structured document content.
    When structured_output is True, expects a DocumentRequest object with specific formatting instructions.
    """
    try:
        logging.info(f"Received request: {request}")
        
        # If this is a structured document request, add specific system instructions
        if request.structured_output and request.document_request:
            system_prompt = f"""You are an AI assistant specialized in generating educational content for K-12 PE, Health, and Driver's Ed.
            You are generating content for a {request.document_request.document_type}.
            Format the output according to the provided structure and maintain professional educational standards."""
            
            # Convert the document request to a format suitable for GPT
            formatted_prompt = f"""
            Title: {request.document_request.title}
            Document Type: {request.document_request.document_type}
            
            Please provide content for each section while maintaining educational best practices and standards.
            
            Original prompt: {request.prompt}
            """
        else:
            system_prompt = "You are an AI assistant that generates text based on prompts."
            formatted_prompt = request.prompt

        try:
            if not OPENAI_API_KEY:
                raise ValueError("OpenAI API key is not configured")

            # Try OpenAI API first
            client = get_openai_client()
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": formatted_prompt}
                ],
                max_tokens=1000  # Add reasonable limit
            )
            generated_content = response.choices[0].message.content.strip()
            logging.info("Successfully generated content using OpenAI API")
        except ValueError as ve:
            logging.error(f"Configuration error: {str(ve)}")
            raise HTTPException(status_code=503, detail="Text generation service is not configured")
        except openai.RateLimitError:
            logging.error("OpenAI API rate limit exceeded")
            raise HTTPException(status_code=429, detail="Rate limit exceeded. Please try again later.")
        except openai.AuthenticationError:
            logging.error("OpenAI API authentication failed")
            raise HTTPException(status_code=503, detail="Text generation service authentication failed")
        except Exception as openai_error:
            logging.error(f"OpenAI API error: {str(openai_error)}")
            
            # Fallback: Generate a basic response based on the document type
            if request.structured_output and request.document_request:
                generated_content = f"""
                # {request.document_request.title}
                
                This is a placeholder content for your {request.document_request.document_type}.
                The OpenAI API is currently unavailable ({str(openai_error)}).
                
                Please try again later or contact support if the issue persists.
                """
            else:
                generated_content = f"""
                Unable to generate content using AI at the moment.
                Error: {str(openai_error)}
                
                Please try again later or contact support if the issue persists.
                """
        
        # If this is a structured request, process the content for document generation
        if request.structured_output and request.document_request:
            # Process the generated content based on document type
            if request.document_request.output_format == "docx":
                doc = Document()
                doc.add_heading(request.document_request.title, 0)
                doc.add_paragraph(generated_content)
                    
                # Save with metadata
                filename = f"{request.document_request.document_type}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
                doc.save(filename)
                
                return FileResponse(filename)
            
            elif request.document_request.output_format == "pptx":
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                title = slide.shapes.title
                title.text = request.document_request.title
                
                # Add content to the slide
                left = top = width = height = Inches(1)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = generated_content
                
                filename = f"{request.document_request.document_type}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                prs.save(filename)
                return FileResponse(filename)
            
        logging.info("Content generation successful")
        return {"generated_text": generated_content, "source": "openai" if "openai_error" not in locals() else "fallback"}
        
    except Exception as e:
        logging.error(f"Error in generate-text: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/test")
async def test_endpoint():
    """Simple endpoint to verify API is accessible"""
    return {"status": "ok", "message": "API is running and accessible"}

@app.post("/generate-document")
async def generate_document(request: DocumentRequest):
    """
    Simplified document generation endpoint.
    Accepts content and generates a document in the requested format.
    """
    try:
        logging.info(f"Generating {request.output_format} document: {request.title}")
        
        if request.output_format == "docx":
            doc = Document()
            doc.add_heading(request.title, 0)
            doc.add_paragraph(request.content)
            
            filename = f"generated_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
            doc.save(filename)
            return FileResponse(filename)
            
        elif request.output_format == "pptx":
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = request.title
            
            # Add content to the slide
            left = top = width = height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = request.content
            
            filename = f"generated_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            prs.save(filename)
            return FileResponse(filename)
            
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported format: {request.output_format}")
            
    except Exception as e:
        logging.error(f"Error generating document: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)

